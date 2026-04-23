#!/usr/bin/env python3
"""Build the DakCUPAC Bourbon Tasting as an actual PowerPoint pitch deck.

Mirrors the content of bourbon-tasting-2026.docx but in .pptx form — real
slides, title + content layout, clickable hyperlinks on each source.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


NAVY = RGBColor(0x1F, 0x2A, 0x44)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0x88, 0x88, 0x88)
LINK = RGBColor(0x1F, 0x4E, 0x79)


def add_title_slide(prs, title, subtitle, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.1), Inches(1.3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = NAVY

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.6))
    p = sub_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    r.font.size = Pt(20)
    r.font.color.rgb = BLUE
    r.font.italic = True

    # Body
    body_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(12.1), Inches(4.0))
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(14)
        r.font.color.rgb = GRAY
    return slide


def add_content_slide(prs, title, items, footer_text):
    """items: list of either str (plain bullet) or dict with keys
    {prefix, link_text, url, suffix} or {bold: True, text: str} or {heading: str}."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = NAVY

    # Accent bar
    from pptx.shapes.autoshape import Shape
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.25), Inches(1.3), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    # Body
    body_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5))
    tf = body_box.text_frame
    tf.word_wrap = True

    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        if isinstance(item, str):
            # Plain bullet
            p.level = 0
            r = p.add_run()
            r.text = "•  " + item
            r.font.size = Pt(14)
            r.font.color.rgb = GRAY
            p.space_after = Pt(6)
        elif item.get("heading"):
            r = p.add_run()
            r.text = item["heading"]
            r.font.size = Pt(16)
            r.font.bold = True
            r.font.color.rgb = BLUE
            p.space_after = Pt(4)
            p.space_before = Pt(8)
        elif item.get("bold"):
            r = p.add_run()
            r.text = item["text"]
            r.font.size = Pt(14)
            r.font.bold = True
            r.font.color.rgb = NAVY
            p.space_after = Pt(6)
        elif item.get("blank"):
            r = p.add_run()
            r.text = ""
            r.font.size = Pt(6)
        elif item.get("link"):
            # "• prefix LINKTEXT suffix"
            r = p.add_run()
            r.text = "•  " + item.get("prefix", "")
            r.font.size = Pt(13)
            r.font.color.rgb = GRAY

            r2 = p.add_run()
            r2.text = item["link_text"]
            r2.font.size = Pt(13)
            r2.font.color.rgb = LINK
            r2.font.underline = True
            r2.hyperlink.address = item["url"]

            if item.get("suffix"):
                r3 = p.add_run()
                r3.text = item["suffix"]
                r3.font.size = Pt(13)
                r3.font.color.rgb = GRAY
            p.space_after = Pt(6)

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.3))
    p = footer_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = footer_text
    r.font.size = Pt(9)
    r.font.italic = True
    r.font.color.rgb = LIGHT
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    TOTAL = 11

    # --- Slide 1: Title ---
    add_title_slide(
        prs,
        "DakCUPAC Bourbon & Whiskey Tasting",
        "Fundraiser pitch deck — subcommittee working draft",
        [
            "Prepared April 23, 2026. All linked costs verified against live listings.",
            "",
            "Event concept: private, instructor-led tasting for CU executives, board members, and CU-affiliated PAC donors. Six pours, 90 minutes, Scott “Fuzzy” Meske leading.",
            "",
            "Engraved Glencairn glass gifted at the door. Chick-fil-A catering.",
            "",
            "Target: 40 attendees × $250/ticket = $10,000 gross.",
            "",
            "Per subcommittee chair: the only costs tracked in this deck are glassware and catering. Bottles are covered by CU sponsors. Venue and instructor fee carried separately.",
        ],
    )

    # --- Slide 2: Why this event works ---
    add_content_slide(
        prs,
        "Why this event works",
        [
            "Right audience for the ticket price. Credit union execs and board members sit exactly where the federal/state PAC contribution ceiling lives. $250 reads as a bourbon-tasting cost, not a fundraiser ask.",
            "Sponsorships eat the biggest cost line. CU bottle sponsors cover the whiskey; PAC keeps the gate.",
            "The glass is the marketing. Every attendee leaves with an engraved Glencairn that lives in their bar. DakCUPAC shows up every time they pour for the next decade.",
            "Scott is a force multiplier. Guided tasting is a different product than “open bar with bourbon.” It justifies the ticket and gives the room a shared experience.",
        ],
        f"Slide 2 of {TOTAL}",
    )

    # --- Slide 3: The numbers ---
    add_content_slide(
        prs,
        "The numbers at a glance",
        [
            {"bold": True, "text": "Gross revenue: 40 × $250 = $10,000"},
            {"blank": True},
            {"heading": "Tracked costs"},
            "Engraved Glencairn glasses (48 at $15.90 each, 48+ tier 15% off): $763.20",
            "Chick-fil-A catering (40 — spicy + non-spicy): $480 – $700",
            {"bold": True, "text": "Tracked-cost subtotal: $1,243 – $1,463"},
            {"bold": True, "text": "Net to PAC: $8,537 – $8,757"},
            {"blank": True},
            {"heading": "Carried separately (not tracked here)"},
            "Bottles — covered by six CU bottle sponsors",
            "Venue — Elks Lodge or private residence (TBD)",
            "Scott’s instructor fee (TBD)",
        ],
        f"Slide 3 of {TOTAL}",
    )

    # --- Slide 4: Glassware ---
    add_content_slide(
        prs,
        "Glassware — 48 engraved Glencairn glasses",
        [
            {"bold": True, "text": "Confirmed tier: 48+ qty at 15% off = $15.90 USD each."},
            "Order 48, not 40 — hits the discount tier; spares cover breakage, late RSVPs, and the sponsoring CUs.",
            {"bold": True, "text": "Math: 48 × $15.90 = $763.20"},
            {"blank": True},
            {"heading": "Vendor options"},
            {
                "link": True,
                "prefix": "",
                "link_text": "Sip n Savor — bulk engraved Glencairn (48+ tier)",
                "url": "https://sipnsavor.co/products/copy-of-bulk-glencairn-whisky-glasses",
                "suffix": " — $15.90/glass at 48+, free shipping over $100",
            },
            {
                "link": True,
                "prefix": "",
                "link_text": "Quality Glass Engraving — Stölzle Glencairn 6oz",
                "url": "https://qualityglassengraving.com/products/engraved-stolzle-glencairn-6-oz-whiskey-glass-item-3550031t",
                "suffix": " — UV laser, no setup fee, 5–10 day turnaround",
            },
            {
                "link": True,
                "prefix": "",
                "link_text": "Crystal Imagery — Glencairn collection",
                "url": "https://crystalimagery.com/collections/glencairn-whisky-glass",
                "suffix": " — deep-carved engraving, quote-based",
            },
            {
                "link": True,
                "prefix": "",
                "link_text": "Distillery Products — wholesale Glencairn",
                "url": "https://www.distilleryproducts.com/wholesale-glencairn/",
                "suffix": " — 48-case blanks at $5.90/glass before engraving",
            },
            {
                "link": True,
                "prefix": "",
                "link_text": "Glencairn Crystal — official wholesale",
                "url": "https://business.glencairn.com/wholesale-whisky-glasses/",
                "suffix": "",
            },
            {"blank": True},
            {"bold": True, "text": "Recommendation: Sip n Savor at the confirmed tier. Artwork at T-10, delivery deadline T-3."},
        ],
        f"Slide 4 of {TOTAL}",
    )

    # --- Slide 5: Catering ---
    add_content_slide(
        prs,
        "Catering — Chick-fil-A for 40 attendees",
        [
            "Per chair: both spicy and non-spicy sandwich trays so attendees can pick.",
            {"blank": True},
            {"heading": "Menu items (official Chick-fil-A catering)"},
            {
                "link": True,
                "prefix": "",
                "link_text": "Large Hot Nuggets Tray",
                "url": "https://www.chick-fil-a.com/catering/catering-trays/large-hot-chick-fil-a-nuggets-tray",
                "suffix": " — 200 nuggets, serves ~25",
            },
            {
                "link": True,
                "prefix": "",
                "link_text": "Sandwich Tray (classic / non-spicy)",
                "url": "https://www.chick-fil-a.com/catering/trays",
                "suffix": "",
            },
            {
                "link": True,
                "prefix": "",
                "link_text": "Spicy Chicken — confirm tray availability with operator",
                "url": "https://www.chick-fil-a.com/catering",
                "suffix": "",
            },
            "Side trays: mac & cheese, waffle fries, fruit tray",
            {"blank": True},
            {"heading": "Bismarck ND specifics"},
            {
                "link": True,
                "prefix": "",
                "link_text": "Chick-fil-A Kirkwood Mall (Bismarck ND)",
                "url": "https://www.chick-fil-a.com/locations/nd/kirkwood-mall-nd",
                "suffix": " — catering delivery minimum $175",
            },
            "Bismarck catering line: (701) 751-0793",
            "Prices are location-gated — call or enter ZIP for a firm Bismarck quote",
            {"blank": True},
            {"bold": True, "text": "Planned order (spicy + non-spicy): ~$436 food + delivery/tip → budget $480 – $700"},
        ],
        f"Slide 5 of {TOTAL}",
    )

    # --- Slide 6: Bottle lineup ---
    add_content_slide(
        prs,
        "Bottle lineup — six pours, buyable around $500",
        [
            "Each bottle links to a real retailer where a CU sponsor can buy it today. Bottle cost is sponsor-covered, not a tracked PAC cost.",
            {"blank": True},
            {
                "link": True,
                "prefix": "",
                "link_text": "Michter’s 10 Year Single Barrel Bourbon (2026) — Cana Wine Co.",
                "url": "https://canawineco.com/products/2026-michters-10-year-old-single-barrel-bourbon-whiskey-750ml",
                "suffix": " — $549.99",
            },
            {
                "link": True,
                "prefix": "",
                "link_text": "Old Forester Birthday Bourbon 2025 — Huntsman Heritage",
                "url": "https://huntsmanheritage.com/products/old-forester-birthday-bourbon-2025",
                "suffix": " — $699.99",
            },
            {
                "link": True,
                "prefix": "",
                "link_text": "The Macallan 18 Year Sherry Oak (2026) — Flaviar",
                "url": "https://flaviar.com/products/the-macallan-18-year-old-sherry-oak-2025-release-single-malt-scotch-whisky-750",
                "suffix": " — ~$412",
            },
            {
                "link": True,
                "prefix": "",
                "link_text": "Redbreast 21 Year Irish Whiskey — Total Wine",
                "url": "https://www.totalwine.com/spirits/irish-whiskey/redbreast-21-yr-irish-whiskey/p/137936750",
                "suffix": " — $400–500",
            },
            {
                "link": True,
                "prefix": "",
                "link_text": "Willett Family Estate 8 Year Single Barrel — Blackwell’s Wines",
                "url": "https://www.blackwellswines.com/products/willett-family-estate-8-year-bourbon-whiskey",
                "suffix": " — $300–550",
            },
            {
                "link": True,
                "prefix": "",
                "link_text": "Parker’s Heritage Collection — The Bourbon Concierge",
                "url": "https://thebourbonconcierge.com/collections/parkers-heritage",
                "suffix": " — $500+ routinely",
            },
            {"blank": True},
            {"heading": "Optional 7th — host pour"},
            {
                "link": True,
                "prefix": "",
                "link_text": "Heart River Spirits — Painted Canyon Bourbon",
                "url": "https://heartriverspirits.com/paintedcanyonbourbon/",
                "suffix": " (Bismarck craft distiller — quote direct)",
            },
        ],
        f"Slide 6 of {TOTAL}",
    )

    # --- Slide 7: Pricing sources ---
    add_content_slide(
        prs,
        "Pricing sources",
        [
            {"heading": "Glassware"},
            {"link": True, "prefix": "", "link_text": "Sip n Savor — bulk Glencairn", "url": "https://sipnsavor.co/products/copy-of-bulk-glencairn-whisky-glasses"},
            {"link": True, "prefix": "", "link_text": "Quality Glass Engraving — Stölzle Glencairn 6oz", "url": "https://qualityglassengraving.com/products/engraved-stolzle-glencairn-6-oz-whiskey-glass-item-3550031t"},
            {"link": True, "prefix": "", "link_text": "Crystal Imagery — Glencairn collection", "url": "https://crystalimagery.com/collections/glencairn-whisky-glass"},
            {"link": True, "prefix": "", "link_text": "Distillery Products wholesale", "url": "https://www.distilleryproducts.com/wholesale-glencairn/"},
            {"link": True, "prefix": "", "link_text": "Glencairn Crystal wholesale", "url": "https://business.glencairn.com/wholesale-whisky-glasses/"},
            {"heading": "Catering"},
            {"link": True, "prefix": "", "link_text": "Chick-fil-A — Large Hot Nuggets Tray", "url": "https://www.chick-fil-a.com/catering/catering-trays/large-hot-chick-fil-a-nuggets-tray"},
            {"link": True, "prefix": "", "link_text": "Chick-fil-A — catering trays overview", "url": "https://www.chick-fil-a.com/catering/trays"},
            {"link": True, "prefix": "", "link_text": "Chick-fil-A — catering landing page", "url": "https://www.chick-fil-a.com/catering"},
            {"link": True, "prefix": "", "link_text": "Chick-fil-A — Kirkwood Mall, Bismarck ND", "url": "https://www.chick-fil-a.com/locations/nd/kirkwood-mall-nd"},
        ],
        f"Slide 7 of {TOTAL}",
    )

    # --- Slide 8: Sponsorship ---
    add_content_slide(
        prs,
        "Sponsorship — six CU bottle sponsors",
        [
            "Each sponsoring CU covers one bottle in the tasting flight (outside tracked PAC costs).",
            {"blank": True},
            {"heading": "Each sponsor receives"},
            "Tent card next to their bottle: “This pour proudly sponsored by [CU + logo]”",
            "Verbal callout from Scott when the pour is introduced",
            "Recognition on the printed tasting menu",
            "Logo on check-in event signage",
            "Two complimentary tickets (worth $500)",
            {"blank": True},
            {"heading": "Optional add-on — Featured ND Pour sponsor"},
            "Underwrites the host pour (Heart River Painted Canyon)",
            "Lower buy-in; good fit for an ND-HQ CU",
            "One complimentary ticket included",
        ],
        f"Slide 8 of {TOTAL}",
    )

    # --- Slide 9: Timeline ---
    add_content_slide(
        prs,
        "Timeline & milestones",
        [
            "Working back from a target Thursday evening, late summer / early fall 2026.",
            {"blank": True},
            "T-12 weeks — confirm date with Scott; lock venue; call Bismarck Chick-fil-A for written quote",
            "T-10 weeks — sponsor outreach to six CUs; finalize bottle list; submit Glencairn artwork",
            "T-8 weeks — sponsorships in; glass deposit paid; venue contract signed",
            "T-6 weeks — open ticket sales; catering order placed (revisable)",
            "T-4 weeks — bottle sourcing confirmed; printed materials to printer",
            "T-3 weeks — glass delivery deadline",
            "T-2 weeks — final headcount; catering spicy/non-spicy split confirmed",
            "T-1 week — venue walkthrough; AV check; glass delivery confirmed",
            "Day-of — setup 4:00 PM; doors 5:30; tasting 6:00",
        ],
        f"Slide 9 of {TOTAL}",
    )

    # --- Slide 10: Risks ---
    add_content_slide(
        prs,
        "Risks & mitigations",
        [
            "Glass delivery slips — T-10 order, hard T-3 delivery date, non-engraved backup vendor on standby",
            "Chick-fil-A spicy tray not offered at Bismarck — confirm at T-12; fallback is classic tray + spicy singles from the main menu",
            "Catering pricing varies in ND — get a written quote from the Bismarck operator at T-12, not T-2",
            "$175 catering delivery minimum — already well cleared",
            "Headcount lands below 40 — couples-pricing or “bring a colleague” upsell at T-4; food and glass costs scale down cleanly",
        ],
        f"Slide 10 of {TOTAL}",
    )

    # --- Slide 11: The ask ---
    add_content_slide(
        prs,
        "The ask",
        [
            {"bold": True, "text": "Subcommittee, decide today:"},
            {"blank": True},
            "Approve the $250 / 40-attendee anchor — $10,000 gross target.",
            "Approve the glassware line: 48 engraved Glencairn glasses at $15.90 each = $763.20.",
            "Approve the catering budget: $480 – $700 for Chick-fil-A (spicy + non-spicy, nuggets, sides).",
            "Authorize the chair to open the conversation with Scott.",
            {"blank": True},
            {"bold": True, "text": "Tracked-cost subtotal: $1,243 – $1,463 against a $10,000 gross."},
            {"bold": True, "text": "Net to PAC: $8,537 – $8,757."},
        ],
        f"Slide 11 of {TOTAL}",
    )

    out_path = "/home/user/dakcupac-2026/events/bourbon-tasting-2026.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    main()
